from pptx import Presentation
import logging

from data_loaders.data_loader_interface import DataLoaderInterface

class PptxLoader(DataLoaderInterface):
    
    def load_data(self, data_path):
        text = self.convert_pptx_to_txt(data_path)
        return text
    
    def convert_pptx_to_txt(self, filepath):
        try:
            pres = Presentation(filepath)
        except Exception as e:
            logging.error(f"Error opening PPTX file {filepath}: {e}")
            return ""

        text = []
        for slide in pres.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text.append(run.text)

        return " ".join(text)